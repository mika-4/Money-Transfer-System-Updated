from pptx import Presentation
from pptx.util import Inches, Pt
import os

base = os.path.dirname(__file__)
img_path = os.path.join(base, 'architecture-diagram.png')
output = os.path.join(base, 'money-transfer-system.pptx')

prs = Presentation()
# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = 'Money Transfer System — Overview'
slide.placeholders[1].text = 'Angular frontend + Spring Boot backend'

# Agenda slide
bullet_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_layout)
slide.shapes.title.text = 'Agenda'
body = slide.shapes.placeholders[1].text_frame
body.text = 'Overview: System purpose and users'
for txt in ['Architecture: Components & data flow', 'Unique Features: Key differentiators', 'Wrap-up: Next steps / Q&A']:
    p = body.add_paragraph()
    p.text = txt
    p.level = 1

# Architecture slide with image
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'Architecture'
left = Inches(0.5)
top = Inches(1.6)
width = Inches(9)
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, left, top, width=width)
else:
    tx = slide.shapes.add_textbox(left, top, width, Inches(1.5))
    tx.text = 'Architecture image not found.'

# Unique features slide
slide = prs.slides.add_slide(bullet_layout)
slide.shapes.title.text = 'Unique Features'
body = slide.shapes.placeholders[1].text_frame
body.text = 'Modular Frontend: Components, services, guards, interceptors'
features = [
    'Layered Backend: Controller → Service → Repository',
    'Token-based Security: interceptor + route guards',
    'Transactional Transfers with Audit Trail',
    'Extensible Integrations: Payment gateway & notifications',
    'User Experience: Transfer confirmation and history view'
]
for f in features:
    p = body.add_paragraph()
    p.text = f
    p.level = 1

# Speaker notes slide
slide = prs.slides.add_slide(bullet_layout)
slide.shapes.title.text = 'Speaker Notes'
body = slide.shapes.placeholders[1].text_frame
body.text = 'Title: Introduce purpose — secure, modular money transfer demo.'
notes = [
    'Agenda: Quick walkthrough of slides and time allocation.',
    'Architecture: Explain data flow from browser to DB and external services; mention interceptor/guard role.',
    'Unique features: Highlight transaction safety, audit/history, and easy integrations.',
    'Wrap-up: Next steps — demo, metrics, or convert slides to PPTX on request.'
]
for n in notes:
    p = body.add_paragraph()
    p.text = n
    p.level = 1

prs.save(output)
print('Saved PPTX to', output)
